import os
import json
import logging
import torch
import time
from typing import List, Dict, Any, Optional, Tuple, Set
from pathlib import Path
from dataclasses import dataclass
import numpy as np
import faiss
from sentence_transformers import SentenceTransformer
import streamlit as st
import ollama
from PIL import Image
import fitz
from docx import Document
from pptx import Presentation
import io
from functools import lru_cache
from concurrent.futures import ThreadPoolExecutor, as_completed
import re
import sys
OLLAMA_BASE_URL = "https://local-search-engine-hzjmswr4j4mkyrys5fxvmc.streamlit.app:11434"
ollama.client = ollama.Client(host=OLLAMA_BASE_URL)
# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Constants
CHUNK_SIZE = 500
CHUNK_OVERLAP = 50
MODEL_NAME = 'sentence-transformers/all-MiniLM-L6-v2'  # Faster, lighter model
INDEX_PATH = "document_index.faiss"
METADATA_PATH = "metadata.json"
MAX_WORKERS = 4  # Adjust based on your CPU cores

@dataclass
class SearchResult:
    id: int
    path: str
    content: str
    score: float

class DocumentProcessor:
    def __init__(self):
        self.model = SentenceTransformer(MODEL_NAME)
        self.dimension = self.model.get_sentence_embedding_dimension()
        self.index = None
        self.metadata: List[Dict[str, Any]] = []
        self._initialize_index()
        self.document_cache = {}  # Cache for document content

    def _initialize_index(self):
        """Initialize or load existing FAISS index"""
        try:
            if os.path.exists(INDEX_PATH) and os.path.exists(METADATA_PATH):
                self.index = faiss.read_index(INDEX_PATH)
                with open(METADATA_PATH, 'r') as f:
                    self.metadata = json.load(f)
                logger.info("Loaded existing index and metadata")
            else:
                self.index = faiss.IndexFlatIP(self.dimension)
                logger.info("Created new index")

        except Exception as e:
            logger.error(f"Error initializing index: {e}")
            self.index = faiss.IndexFlatIP(self.dimension)

    @lru_cache(maxsize=1000)
    def chunk_text(self, text: str) -> List[str]:
        """Split text into manageable chunks with caching"""
        words = text.split()
        chunks = []
        for i in range(0, len(words), CHUNK_SIZE - CHUNK_OVERLAP):
            chunk = ' '.join(words[i:i + CHUNK_SIZE])
            chunks.append(chunk)
        return chunks 

    def process_file(self, file_path: str) -> tuple:
        """Process a single file and return its chunks and metadata"""
        content = self.read_file(str(file_path))
        if not content:
            return [], []
        chunks = self.chunk_text(content)
        file_metadata = [{"path": str(file_path), "chunk_id": i} for i in range(len(chunks))]
        return chunks, file_metadata

    def read_file(self, file_path: str) -> str:
        """Read content from different file types with caching"""
        if file_path in self.document_cache:
            return self.document_cache[file_path]
        try:
            ext = Path(file_path).suffix.lower()
            content = ""
            if ext == '.txt':
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            elif ext == '.pdf':
                with fitz.open(file_path) as doc:
                    content = " ".join(page.get_text() for page in doc)
            elif ext == '.docx':
                doc = Document(file_path)
                content = " ".join(p.text for p in doc.paragraphs)
            elif ext == '.pptx':
                prs = Presentation(file_path)
                content = " ".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
            self.document_cache[file_path] = content
            return content
        except Exception as e:
            logger.error(f"Error reading file {file_path}: {e}")
            return ""

    def index_documents(self, directory: str) -> bool:
        """Index documents from directory using parallel processing"""
        try:
            file_paths = [p for ext in ['.txt', '.pdf', '.docx', '.pptx'] for p in Path(directory).rglob(f'*{ext}')]

            if not file_paths:
                st.warning("No supported documents found")
                return False

            documents = []
            self.metadata = []
            with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
                future_to_file = {executor.submit(self.process_file, str(file_path)): file_path for file_path in file_paths}
                for future in as_completed(future_to_file):
                    chunks, file_metadata = future.result()
                    documents.extend(chunks)
                    self.metadata.extend(file_metadata)

            if not documents:
                st.warning("No content extracted from documents")
                return False

            batch_size = 32
            all_embeddings = [self.model.encode(documents[i:i + batch_size], show_progress_bar=False, batch_size=batch_size) for i in range(0, len(documents), batch_size)]
            embeddings = np.vstack(all_embeddings)

            self.index = faiss.IndexFlatIP(self.dimension)
            self.index.add(embeddings.astype('float32'))

            faiss.write_index(self.index, INDEX_PATH)
            with open(METADATA_PATH, 'w') as f:
                json.dump(self.metadata, f)

            st.success(f"Indexed {len(documents)} text chunks from {len(file_paths)} files")
            return True

        except Exception as e:
            logger.error(f"Indexing error: {e}")
            st.error(f"Error during indexing: {str(e)}")
            return False

    @lru_cache(maxsize=100)
    def search(self, query: str, k: int = 5) -> List[SearchResult]:
        """Search for relevant documents with caching"""
        try:
            if not self.index or self.index.ntotal == 0:
                st.warning("No documents indexed yet")
                return []

            # Encode query and search
            query_vector = self.model.encode([query])[0]
            distances, indices = self.index.search(
                query_vector.reshape(1, -1).astype('float32'),
                min(k, self.index.ntotal)
            )

            # Process results
            results = []
            for i, idx in enumerate(indices[0]):
                if idx < len(self.metadata):
                    meta = self.metadata[idx]
                    try:
                        content = self.read_file(meta["path"])
                        chunks = self.chunk_text(content)
                        chunk_content = chunks[meta["chunk_id"]] if meta["chunk_id"] < len(chunks) else ""

                        results.append(SearchResult(
                            id=int(idx),
                            path=meta["path"],
                            content=chunk_content,
                            score=float(distances[0][i])
                        ))

                    except Exception as e:
                        logger.error(f"Error reading chunk: {e}")
                        continue

            return results
        except Exception as e:
            logger.error(f"Search error: {e}")
            st.error(f"Error during search: {str(e)}")
            return []
@lru_cache(maxsize=100)
def generate_answer(query: str, context: str) -> Tuple[str, Set[int]]:
    """
    Generate answer using Ollama with caching and return referenced citations
    Args:

        query (str): The user's question
        context (str): The context information from documents

    Returns:

        Tuple[str, Set[int]]: A tuple containing the generated answer and a set of citation indices
    """
    try:

        prompt = f"""
        Answer based on the context below. Be concise.
        You must cite your sources using [0], [1], etc. for EVERY claim you make.
        Make sure to use the citations explicitly in your answer.
        Context:

        {context}

        Question: {query}

        Answer:"""
        response = ollama.generate(
            model='qwen2.5-coder:14b',
            prompt=prompt
        )

        answer = response['response']
        logger.info(f"answer: {answer}")
        # Extract citation numbers from the answer
        citations = set(int(num) for num in re.findall(r'\[(\d+)\]', answer))
        logger.info(f"citations: {citations}")
        return answer, citations

    except Exception as e:
        logger.error(f"Answer generation error: {e}")
        return "Sorry, I couldn't generate an answer at this time.", set()
def main():
    st.set_page_config(page_title="Local Search", layout="wide")
    
    if 'processor' not in st.session_state:
        st.session_state.processor = DocumentProcessor()

    st.title("Local Document Search")

    # docs_path = st.text_input("Documents folder path:")
    # if st.button("Index Documents") and docs_path:
        # with st.spinner("Indexing documents..."):
    st.session_state.processor.index_documents('files')

    query = st.text_input("Enter your question:")
    if st.button("Search") and query:
        with st.spinner("Searching..."):
            results = st.session_state.processor.search(query)
            logger.info(f"results: {results}")
            if results:
                context = "\n\n".join(f"[{i}] {r.content}" for i, r in enumerate(results))
                answer, citations = generate_answer(query, context)
                st.subheader("Answer")
                st.write(answer)

                if citations:
                    st.subheader("Referenced Documents")
                    for citation_num in sorted(citations):
                        if citation_num < len(results):
                            result = results[citation_num]
                            with st.expander(f"Reference [{citation_num}] - {Path(result.path).name}") :
                                st.write(result.content)
                else:
                    st.info("No specific documents were cited in the answer.")
            else:
                st.warning("No relevant documents found")

if __name__ == "__main__":
    main()
