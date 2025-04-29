import os
import json
import logging
import torch
import time
from typing import List, Dict, Any, Optional, Tuple, Set
from pathlib import Path
from dataclasses import dataclass
import numpy as np
import faiss
from sentence_transformers import SentenceTransformer
import streamlit as st
import ollama
from PIL import Image
import fitz
from docx import Document
from pptx import Presentation
import io
from functools import lru_cache
from concurrent.futures import ThreadPoolExecutor, as_completed
import re
import sys
import mimetypes
import requests

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)
ollama.BASE_URL = "http://ollama:11434"

# Constants
CHUNK_SIZE = 500
CHUNK_OVERLAP = 50
MODEL_NAME = 'sentence-transformers/all-MiniLM-L6-v2'  # Faster, lighter model
INDEX_PATH = "document_index.faiss"
METADATA_PATH = "metadata.json"
MAX_WORKERS = 4  # Adjust based on your CPU cores

@dataclass
class SearchResult:
    id: int
    path: str
    content: str
    score: float

class DocumentProcessor:
    def __init__(self):
        self.model = SentenceTransformer(MODEL_NAME)
        self.dimension = self.model.get_sentence_embedding_dimension()
        # Initialize FAISS index inside __init__
        if hasattr(faiss, "IndexHNSWFlat"):
            self.index = faiss.IndexHNSWFlat(self.dimension, 32)
        else:
            self.index = faiss.IndexFlatL2(self.dimension)  # Fallback for unsupported FAISS version
        self.index = None
        self.metadata: List[Dict[str, Any]] = []
        self._initialize_index()
        self.document_cache = {}  # Cache for document content

    def _initialize_index(self):
        """Initialize or load FAISS index"""
        try:
            if os.path.exists(INDEX_PATH) and os.path.exists(METADATA_PATH):
                with open(METADATA_PATH, 'r') as f:
                    self.metadata = json.load(f)
                logger.info("Loaded existing FAISS index")
            else:
                logger.info("Created new HNSW index")
        except Exception as e:
            logger.error(f"Error initializing index: {e}")

    def read_file_binary(self, file_path: str) -> bytes:
        try:
            with open(file_path, 'rb') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Error reading file {file_path} in binary mode: {e}")
            return b""

    @lru_cache(maxsize=1000)
    def chunk_text(self, text: str) -> List[str]:
        """Split text into manageable chunks with caching"""
        words = text.split()
        chunks = []
        for i in range(0, len(words), CHUNK_SIZE - CHUNK_OVERLAP):
            chunk = ' '.join(words[i:i + CHUNK_SIZE])
            chunks.append(chunk)
        return chunks 

    def process_file(self, file_path: str) -> tuple:
        """Process a single file and return its chunks and metadata"""
        content = self.read_file(str(file_path))
        if not content:
            return [], []
        chunks = self.chunk_text(content)
        file_metadata = [{"path": str(file_path), "chunk_id": i} for i in range(len(chunks))]
        return chunks, file_metadata

    def read_file(self, file_path: str) -> str:
        """Read content from different file types with caching"""
        if file_path in self.document_cache:
            return self.document_cache[file_path]
        try:
            ext = Path(file_path).suffix.lower()
            content = ""
            if ext == '.txt':
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            elif ext == '.pdf':
                with fitz.open(file_path) as doc:
                    content = " ".join(page.get_text() for page in doc)
            elif ext == '.docx':
                doc = Document(file_path)
                content = " ".join(p.text for p in doc.paragraphs)
            elif ext == '.pptx':
                prs = Presentation(file_path)
                content = " ".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
            self.document_cache[file_path] = content
            return content
        except Exception as e:
            logger.error(f"Error reading file {file_path}: {e}")
            return ""

    def index_documents(self, directory: str) -> bool:
        """Index documents from directory using parallel processing"""
        try:
            file_paths = [p for ext in ['.txt', '.pdf', '.docx', '.pptx'] for p in Path(directory).rglob(f'*{ext}')]

            if not file_paths:
                st.warning("No supported documents found")
                logger.info("No supported documents found in the directory")
                return False

            logger.info(f"Found {len(file_paths)} documents to index.")
            documents = []
            self.metadata = []
            with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
                future_to_file = {executor.submit(self.process_file, str(file_path)): file_path for file_path in file_paths}
                for future in as_completed(future_to_file):
                    chunks, file_metadata = future.result()
                    documents.extend(chunks)
                    self.metadata.extend(file_metadata)

            if not documents:
                st.warning("No content extracted from documents")
                return False

            batch_size = 32
            all_embeddings = [self.model.encode(documents[i:i + batch_size], show_progress_bar=False, batch_size=batch_size) for i in range(0, len(documents), batch_size)]
            embeddings = np.vstack(all_embeddings)

            self.index = faiss.IndexFlatIP(self.dimension)
            self.index.add(embeddings.astype('float32'))

            faiss.write_index(self.index, INDEX_PATH)
            with open(METADATA_PATH, 'w') as f:
                json.dump(self.metadata, f)

            # st.success(f"Indexed {len(documents)} text chunks from {len(file_paths)} files")
            return True

        except Exception as e:
            logger.error(f"Indexing error: {e}")
            st.error(f"Error during indexing: {str(e)}")
            return False

    @lru_cache(maxsize=100)
    def search(self, query: str, k: int = 5) -> List[SearchResult]:
        """Search for relevant documents with caching"""
        try:
            if not self.index or self.index.ntotal == 0:
                st.warning("No documents indexed yet")
                return []

            # Encode query and search
            query_vector = self.model.encode([query])[0]
            distances, indices = self.index.search(
                query_vector.reshape(1, -1).astype('float32'),
                min(k, self.index.ntotal)
            )

            # Process results
            results = []
            for i, idx in enumerate(indices[0]):
                if idx < len(self.metadata):
                    meta = self.metadata[idx]
                    try:
                        content = self.read_file(meta["path"])
                        chunks = self.chunk_text(content)
                        chunk_content = chunks[meta["chunk_id"]] if meta["chunk_id"] < len(chunks) else ""

                        results.append(SearchResult(
                            id=int(idx),
                            path=meta["path"],
                            content=chunk_content,
                            score=float(distances[0][i])
                        ))
                    except Exception as e:
                        logger.error(f"Error reading chunk: {e}")
                        continue

            return results
        except Exception as e:
            logger.error(f"Search error: {e}")
            st.error(f"Error during search: {str(e)}")
            return []

@lru_cache(maxsize=100)
def generate_answer(query: str, context: str) -> Tuple[str, Set[int]]:
    """
    Generate answer using Ollama with streaming support and return referenced citations
    """
    try:
        prompt = f"""
        Get only the related files based on the question and the provided context.
        Be concise and do not include any external content, references, or URLs.
        Cite sources explicitly from the given context using [0], [1], etc., for every claim made. Only use citations from the provided context.
        Make sure to use the citations explicitly in your answer."
        Context: {context}
        Question: {query}
        Answer:"""
        logger.info(f"prompt: {prompt}")
        start_time = time.time()
        OLLAMA_URL = "http://ollama:11434/api/generate"
        data = {
            "model": "llama3.2:3b",
            "prompt": prompt,
            "stream": True  # Enable streaming
        }

        response = requests.post(OLLAMA_URL, json=data, stream=True)
        end_time = time.time()
        elapsed_time = end_time - start_time
        logger.info(f"Response time: {elapsed_time:.4f} seconds")
        answer = ""
        citations = set()

        # Stream the response in real-time
        for chunk in response.iter_lines():
            if chunk:
                decoded_chunk = json.loads(chunk.decode('utf-8'))
                token = decoded_chunk.get("response", "")
                answer += token

                # Extract citations on the fly
                citations.update(set(int(num) for num in re.findall(r'\[(\d+)\]', answer)))

                # Streamlit dynamic update
                yield answer, citations

    except Exception as e:
        logger.error(f"Answer generation error: {e}")
        yield "Sorry, I couldn't generate an answer at this time.", set()

def main():
    st.set_page_config(page_title="Local Search", layout="wide")
    
    if 'processor' not in st.session_state:
        st.session_state.processor = DocumentProcessor()

    st.title("Local Document Search")

    st.session_state.processor.index_documents('/files')

    query = st.text_input("Enter your question:")
    
    if st.button("Search") and query:
        with st.spinner("Searching..."):
            results = st.session_state.processor.search(query)
            if results:
                context = "\n\n".join(f"[{i}] {r.content}" for i, r in enumerate(results))
                
                st.subheader("Answer")
                answer_placeholder = st.empty()  # Placeholder for streaming response
                
                full_answer = ""
                citations = set()

                for partial_answer, partial_citations in generate_answer(query, context):
                    full_answer = partial_answer
                    citations = partial_citations
                    answer_placeholder.write(full_answer)  # Update UI dynamically
                logger.info(f"citations: {citations}")
                if citations:
                    st.subheader("Referenced Documents")
                    i = 0
                    for citation_num in sorted(citations):
                        if citation_num < len(results):
                            result = results[citation_num]
                            with st.expander(f"Reference [{citation_num}] - {Path(result.path).name}"):
                                file_data = st.session_state.processor.read_file_binary(str(result.path))
                                i = i + 1
                                st.download_button(
                                    label="Download",
                                    data=file_data,
                                    file_name=Path(result.path).name,
                                    mime=mimetypes.guess_type(str(result.path))[0] or "text/plain",
                                    key=f"download_button_{i}"  # Unique key for each button
                                )
                                st.write(result.content)
                else:
                    st.info("No specific documents were cited in the answer.")
            else:
                st.warning("No relevant documents found")

if __name__ == "__main__":
    main()
